import pandas as pd
import numpy as np

df.head() #6 primeras columnas del df
df.tail() #6 ultimas columnas del df
df.info() #resumen del tipo de variables
df.shape #filas  x columnas del df
df.columns #columnas del df
df.describe() #resumen estadístico de las variables numéricas
del df #borra una variable o df
df['Citas'] = df['Citas'].astype('float64') #cambia el tipo de dato
df['Citas'].dtype #tipo de dato

#--------------------------SELECCIONAR--------------------------
df["Citas"] #seleccionar una variable
df.Citas #seleccionar una variable
df[["Citas", "Universidad"]] #seleccionar dos o más variables
df.loc[:, 'Nombre':'Citas'] #selecciona las variables de nombre a citas
df.loc[:, ['Nombre', 'Apellido', 'Citas']] #selecciona las variables especificadas
df.iloc[:, 1:10] #selecciona las variables de la 1 a la 10
df.iloc[:, [1, 10, 4]] #selecciona las variables por posición
df.drop(['Nombre', 'Citas'], axis=1) #eliminar variables 
df.select_dtypes(include='number') #selecciona numéricas 'object' para cadenas

#--------------------------FILTRAR--------------------------
df[df['grado'] == 'Doctorado'] #filtra donde grado sea igual a Doctorado
df.loc[df['Universidad'] == 'CIMAT'] #filtra donde universidad es CIMAT
df.loc[df['Universidad'] == 'CIMAT', ["Nombre", "Universidad", "Residencia"]] #y columnas deseadas
df.query("Universidad == 'CIMAT'") #filtra donde universidad es CIMAT
df.query("Universidad == 'CIMAT'")[["Nombre", "Universidad", "Residencia"]] #y columnas deseadas
df.query("Citas > 500") #filtra a quienes tengan mas de 500 citas
df.query("Universidad == 'CIMAT' & Citas > 50") #dos o más condiciones 	
df.query("Universidad == 'CIMAT' | Residencia  == 'Brasil'") #condición con O


#------------------------AGRUPAR----------------------------
df.groupby("universidad").count() #recuento por universidad
df.groupby("universidad").size().reset_index(name='recuento') # recuento por universidad solo dos variables como resultado

#conteo por Género
conteo = df["Género"].value_counts()
#convertir el resultado en un df
serie_counts = pd.Series(conteo)
df = pd.DataFrame({'Genero': serie_counts.index, 'counts': serie_counts.values})

df.groupby("universidad").agg({ #maximo, media de citas y cuenta por universidad
    "Citas" : ["max", "mean"],
    "universidad": "count"
    })
    
df.groupby(['universidad', 'Género']).size().reset_index(name='Count') #genero por universidad

df.groupby('universidad').sum('Citas') #total de citas por universidades

#maximo y minimo por grupo
max_salario_indices = df.groupby('Género')['Salario'].idxmax()
df.loc[max_salario_indices]

min_salario_indices = df.groupby('Género')['Salario'].idxmin()
df.loc[min_salario_indices]

#------------------------RENOMBRAR COLUMNAS----------------------------
df.rename(columns={"UTrabajo":"Universidad"})
df.rename(columns={'Columna1': 'País', 'Columna2': 'Género', 'Columna3': 'Universidad'}, inplace=True) #renombrar varias columnas

#------------------------COLUMNAS A ETIQUETAS DE FILA----------------------------
df.set_index('Nombre')

#------------------------NUEVA COLUMNA----------------------------
df['NuevaColumna'] = df['Citas']*2 #Sobreescribe el df original
df.assign(Articulos_1=df['Articulos'] / 2) #no sobreescibre el df a menos que se asigne


#****************Crear una nueva variable de acuerdo a condiciones de otra**********
# Definir condiciones y resultados
condiciones = [
    df['Citas'] < 3000,
    (df['Citas'] >= 3000) & (df['Citas'] <= 7000),
    df['Citas'] > 7000
]

resultados = ['Bajo', 'Medio', 'Alto']

# Crear una nueva columna 'NivelCitas' utilizando np.select
df['NivelCitas'] = np.select(condiciones, resultados, default='Desconocido')

#si la columna grado tiene un nombramiento de mujer se agrega femenino si es de hombre masculino, otro NA
condiciones = [
    df['NOBILIS'].str.contains("^DRA.|^MTRA."),
    df['NOBILIS'].str.contains("^DR.|^MTRO.|^LIC")
]

resultados = ['Femenino', 'Masculino']

# Aplicar las condiciones utilizando np.select
df['Género'] = np.select(condiciones, resultados, default=np.nan)

#---------------------------------EDITAR CELDAS-------------------------------
df.at[0, 'Nombre'] = "Andres A." #cambia el nombre por "Andres A."

#----------------------------SEPARAR/UNIR COLUMNAS------------------------------
df[['Area1', 'Area2']] = df['Areas'].str.split(expand=True, n=1)
#unir 
df['NombreCompleto'] = df['Nombre'].str.cat(df['PApellido'], sep=' ')


#----------------------------JOINS------------------------------
data_personas = {'nombre': ['Alice', 'Bob', 'Charlie', 'David'],
                  'trabajo': ['Ingeniera', 'Doctor', 'Profesor', 'Estudiante']}
df_personas = pd.DataFrame(data_personas)

data_universidades = {'nombre': ['Alice', 'Bob', 'Eve'],
                      'universidad': ['MIT', 'Harvard', 'Stanford']}
df_universidades = pd.DataFrame(data_universidades)



pd.merge(df_personas, df_universidades, on='nombre', how='left')
pd.merge(df_personas, df_universidades, on='nombre', how='right')
pd.merge(df_personas, df_universidades, on='nombre', how='outer') #full join
pd.merge(df_universidades, df_personas, on='nombre', how='left', indicator=True).query('_merge == "left_only"').drop('_merge', axis=1) #anti join

#---------------------PASAR A DATOS LONGITUDINALES------------------------
data = {'Pais': ['Argentina', 'Brasil', 'Chile'],
        'Año': [2020, 2020, 2020],
        'Ventas_Ene': [100, 150, 120],
        'Ventas_Feb': [120, 160, 130],
        'Ventas_Mar': [110, 140, 125]}

df = pd.DataFrame(data)
pd.melt(df, id_vars=['Pais', 'Año'], var_name='Mes', value_name='Ventas')

#Para revertir
df_long.pivot_table(index=['Pais', 'Año'], columns='Mes', values='Ventas', aggfunc='first').reset_index()

#---------------------------EXPORTAR DATAFRAMES------------------------
ruta_csv = "/ruta/base.csv" #csv
df.to_csv(ruta_csv, index=False) 

#exportar a excel
!pip install pandas openpyxl

ruta_excel = "/home/jorgemartinez/Descargas/dfexcel.xlsx"
df.to_excel(ruta_excel, index=False, engine='openpyxl')

#------------------------Exportar a POWERPOINT------------------------
pip install matplotlib python-pptx

import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

df.groupby("universidad")[['Citas', 'Articulos']].sum().plot()

image_stream = BytesIO()
plt.savefig(image_stream, format='png')
plt.close()

# Crear una presentación
presentation = Presentation()

# Añadir una diapositiva
slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Diseño de título y contenido

# Agregar la imagen a la diapositiva
left = top = Inches(1)
pic = slide.shapes.add_picture(image_stream, left, top, height=Inches(2))

# Guardar la presentación
presentation.save("nombredelarchivo.pptx")

#---------------------------TABLA LATEX------------------------
!pip install tabulate
from tabulate import tabulate


# Imprimir la tabla en formato LaTeX
print(tabulate(df, tablefmt='latex_raw', headers='keys'))

#---------------------------TABLA HTML------------------------
# Generar el código HTML
codigo_html = "<table>\n"
codigo_html += "<thead>\n<tr>\n<th>Nombre</th>\n<th>Género</th>\n<th>Area</th>\n<th>Residencia</th>\n<th>Universidad</th>\n<th>Perfil</th>\n</tr>\n</thead>\n"
codigo_html += "<tbody>\n"

for i, row in df.iterrows():
    codigo_html += "<tr>\n"
    codigo_html += f"<td>{row['Nombre']}</td>\n"
    codigo_html += f"<td>{row['Género']}</td>\n"
    codigo_html += f"<td>{row['Area']}</td>\n"
    codigo_html += f"<td>{row['Residencia']}</td>\n"
    codigo_html += f"<td>{row['Universidad']}</td>\n"

    if pd.notna(row['Perfil']):
        codigo_html += f"<td><a href=\"{row['Perfil']}\" target=\"_blank\">{row['Perfil']}</a></td>\n"
    else:
        codigo_html += "<td></td>\n"

    codigo_html += "</tr>\n"

codigo_html += "</tbody>\n</table>"

# Guardar el código HTML en un archivo
with open("tablapyhton.html", "w") as f:
    f.write(codigo_html)

